import fitz  # PyMuPDF for PDF extraction
from pptx import Presentation
from openai import OpenAI
import json
import os
import sys
import numpy as np
import faiss
from sentence_transformers import SentenceTransformer
from typing import List, Dict, Any, Optional
import logging
import pickle
from pathlib import Path
import time

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class PromptManager:
    def __init__(self, config_file="prompts.json"):
        self.config_file = config_file
        self.prompts = self._load_prompts()
    
    def _load_prompts(self):
        """Load prompts from JSON configuration file"""
        try:
            with open(self.config_file, 'r', encoding='utf-8') as f:
                return json.load(f)
        except FileNotFoundError:
            raise FileNotFoundError(f"Prompt configuration file '{self.config_file}' not found")
        except json.JSONDecodeError as e:
            raise ValueError(f"Invalid JSON in prompt configuration: {e}")
    
    def get_prompt_config(self, prompt_type="investment_analysis"):
        """Get prompt configuration for specified type"""
        if prompt_type not in self.prompts:
            raise ValueError(f"Prompt type '{prompt_type}' not found in configuration")
        return self.prompts[prompt_type]

class InvestmentKnowledgeBase:
    """Investment knowledge base for RAG"""
    
    @staticmethod
    def get_knowledge_items():
        return [
            {
                "id": 0,
                "topic": "startup_valuation_methods",
                "category": "valuation",
                "content": "Startup valuation methods include revenue multiples (5-15x ARR for SaaS, 2-5x for e-commerce), Discounted Cash Flow with high discount rates (15-25%), Comparable company analysis using public/private comps, Risk-adjusted NPV for early-stage ventures. Pre-revenue startups valued on TAM, team pedigree, and competitive moats. Series A typically valued at 10-20x ARR with strong growth.",
                "tags": ["valuation", "metrics", "series_a"]
            },
            {
                "id": 1,
                "topic": "market_sizing_tam_sam_som",
                "category": "market_analysis", 
                "content": "Market sizing framework: Total Addressable Market (TAM) - global demand for solution category, Serviceable Addressable Market (SAM) - realistic market segment company can serve with current business model, Serviceable Obtainable Market (SOM) - achievable market share based on competitive positioning. Bottom-up analysis preferred: unit economics × addressable customers × market penetration rates.",
                "tags": ["market_sizing", "tam", "framework"]
            },
            {
                "id": 2,
                "topic": "product_market_fit_indicators",
                "category": "traction",
                "content": "Product-Market Fit signals: >40% users 'very disappointed' without product (Sean Ellis test), organic growth >20% month-over-month, retention cohorts flattening after month 6, Net Promoter Score >50, word-of-mouth growth coefficient >0.15. Qualitative indicators: customers pulling product into organization, unsolicited feature requests, competitive win rates >60%, sales cycle compression.",
                "tags": ["pmf", "retention", "growth", "early_stage"]
            },
            {
                "id": 3,
                "topic": "saas_key_metrics_benchmarks",
                "category": "metrics",
                "content": "SaaS metrics benchmarks: LTV:CAC ratio >3:1 (excellent >5:1), CAC payback period <12 months (best <6 months), Monthly churn rate <5% (SMB) or <2% (Enterprise), Net Revenue Retention >110% (best-in-class >130%), Gross margin >70% (SaaS standard >80%), Annual Contract Value growth, Magic Number >1.0 for efficient growth.",
                "tags": ["saas", "metrics", "benchmarks", "ltv_cac"]
            },
            {
                "id": 4,
                "topic": "early_stage_funding_milestones",
                "category": "funding",
                "content": "Early-stage funding milestones: Pre-seed ($100K-$1M) - MVP and initial user feedback, Seed ($500K-$5M) - product-market fit signals and early traction, Series A ($2M-$15M) - proven PMF with $1M+ ARR and scalable go-to-market. Each round should provide 18-24 months runway. Key metrics progression: Pre-seed focuses on engagement, Seed on retention/growth, Series A on unit economics and scalability.",
                "tags": ["funding", "milestones", "early_stage", "seed", "series_a"]
            },
            {
                "id": 5,
                "topic": "team_assessment_framework",
                "category": "team",
                "content": "Founder assessment criteria: Domain expertise and founder-market fit, Previous startup experience (especially successful exits), Technical/business complementarity, Ability to attract top talent, Leadership under pressure, Coachability and learning agility, Fair equity distribution, Strong advisory board. Red flags: founder conflicts, unrealistic expectations, poor communication skills, inability to hire senior talent.",
                "tags": ["founders", "team", "assessment", "red_flags"]
            },
            {
                "id": 6,
                "topic": "competitive_analysis_moats",
                "category": "strategy",
                "content": "Competitive moats for startups: Network effects (value increases with users - strongest moat), Switching costs (data/integration lock-in), Economies of scale, Proprietary data/algorithms, Brand recognition, Regulatory barriers, Exclusive partnerships. Evaluate moat timeline and defensibility. Technology alone rarely provides lasting advantage without network effects or data moats.",
                "tags": ["competitive_advantage", "moats", "defensibility", "network_effects"]
            },
            {
                "id": 7,
                "topic": "unit_economics_analysis",
                "category": "financials",
                "content": "Unit economics fundamentals: Customer Acquisition Cost (CAC) including fully-loaded sales/marketing costs, Customer Lifetime Value (LTV) with realistic churn assumptions, Contribution margin after variable costs, Payback period calculation, Cohort-based analysis for accuracy. Key ratios: LTV:CAC >3:1, Payback <12 months, positive unit economics at scale with improving trends.",
                "tags": ["unit_economics", "cac", "ltv", "cohort_analysis"]
            },
            {
                "id": 8,
                "topic": "growth_strategy_channels",
                "category": "growth",
                "content": "Growth channel evaluation: Product-led growth (PLG) for viral/self-serve products, Sales-led growth for enterprise/complex products, Marketing-led growth for consumer/SMB markets. Channel metrics: Customer acquisition cost by channel, conversion rates, time to value, retention by acquisition source. Successful startups typically master 1-2 primary channels before expanding.",
                "tags": ["growth_strategy", "plg", "sales_led", "channels"]
            },
            {
                "id": 9,
                "topic": "investment_red_flags",
                "category": "risks",
                "content": "Major investment red flags: Unrealistic financial projections without supporting data, weak founding team lacking domain expertise, unclear go-to-market strategy, poor unit economics with no path to profitability, highly competitive market with no differentiation, legal/IP issues, high customer concentration (>50% revenue from single customer), excessive burn rate without proportional growth.",
                "tags": ["red_flags", "risks", "due_diligence", "warning_signs"]
            }
        ]

class FAISSRAGSystem:
    """FAISS-powered RAG system for investment analysis"""
    
    def __init__(self, embedding_model="all-MiniLM-L6-v2"):
        self.embedding_model = SentenceTransformer(embedding_model)
        self.knowledge_base = InvestmentKnowledgeBase.get_knowledge_items()
        self.index = None
        self.embeddings = None
        self.index_path = Path("./faiss_indexes")
        self.index_path.mkdir(exist_ok=True)
        self._build_or_load_index()
    
    def _build_or_load_index(self):
        """Build or load existing FAISS index"""
        index_file = self.index_path / "investment_kb.index"
        metadata_file = self.index_path / "investment_kb_metadata.pkl"
        
        if self._load_existing_index(index_file, metadata_file):
            logger.info("Loaded existing FAISS index")
            return
        
        logger.info("Building new FAISS index...")
        self._build_new_index()
        self._save_index(index_file, metadata_file)
        logger.info(f"Built and saved FAISS index with {len(self.knowledge_base)} documents")
    
    def _load_existing_index(self, index_file, metadata_file):
        """Load existing FAISS index if available"""
        try:
            if index_file.exists() and metadata_file.exists():
                self.index = faiss.read_index(str(index_file))
                with open(metadata_file, 'rb') as f:
                    data = pickle.load(f)
                    self.embeddings = data['embeddings']
                    # Verify knowledge base hasn't changed
                    if len(data['knowledge_base']) == len(self.knowledge_base):
                        return True
        except Exception as e:
            logger.warning(f"Failed to load existing index: {e}")
        return False
    
    def _build_new_index(self):
        """Build new FAISS index from scratch"""
        # Generate embeddings for all knowledge items
        texts = [item["content"] for item in self.knowledge_base]
        self.embeddings = self.embedding_model.encode(
            texts,
            convert_to_numpy=True,
            show_progress_bar=True
        )
        
        # Create FAISS index (Inner Product for cosine similarity)
        dimension = self.embeddings.shape[1]
        self.index = faiss.IndexFlatIP(dimension)
        
        # Normalize embeddings for cosine similarity
        faiss.normalize_L2(self.embeddings)
        
        # Add embeddings to index
        self.index.add(self.embeddings.astype('float32'))
    
    def _save_index(self, index_file, metadata_file):
        """Save FAISS index and metadata"""
        try:
            faiss.write_index(self.index, str(index_file))
            with open(metadata_file, 'wb') as f:
                pickle.dump({
                    'embeddings': self.embeddings,
                    'knowledge_base': self.knowledge_base
                }, f)
        except Exception as e:
            logger.error(f"Failed to save index: {e}")
    
    def retrieve_knowledge(self, query: str, top_k: int = 5, similarity_threshold: float = 0.3):
        """Retrieve relevant knowledge using FAISS search"""
        # Generate query embedding
        query_embedding = self.embedding_model.encode([query], convert_to_numpy=True)
        faiss.normalize_L2(query_embedding)
        
        # Search FAISS index
        scores, indices = self.index.search(query_embedding.astype('float32'), top_k)
        
        # Filter results and build response
        relevant_knowledge = []
        for score, idx in zip(scores[0], indices[0]):
            if score >= similarity_threshold:
                item = self.knowledge_base[idx]
                relevant_knowledge.append({
                    "content": item["content"],
                    "topic": item["topic"],
                    "category": item["category"],
                    "similarity_score": float(score),
                    "tags": item["tags"]
                })
        
        logger.info(f"Retrieved {len(relevant_knowledge)} relevant knowledge items")
        return relevant_knowledge

# 1. Extract text from PDF
def extract_text_from_pdf(pdf_path):
    doc = fitz.open(pdf_path)
    text = "\n".join(page.get_text() for page in doc)
    doc.close()
    return text

# 2. Extract text from PPTX
def extract_text_from_pptx(pptx_path):
    ppt = Presentation(pptx_path)
    text_runs = []
    for slide_num, slide in enumerate(ppt.slides):
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_runs.append(f"[Slide {slide_num + 1}] {shape.text}")
    return "\n".join(text_runs)

# 3. Enhanced analyze function with RAG
def analyze_pitch_deck_with_rag(content, openai_api_key, prompt_manager=None, rag_system=None):
    if prompt_manager is None:
        prompt_manager = PromptManager()
    
    if rag_system is None:
        rag_system = FAISSRAGSystem()
    
    # Retrieve relevant knowledge using RAG
    logger.info("Retrieving relevant investment knowledge...")
    relevant_knowledge = rag_system.retrieve_knowledge(content, top_k=7)
    
    # Build context from retrieved knowledge
    context_sections = []
    for item in relevant_knowledge:
        context_sections.append(
            f"**{item['topic'].replace('_', ' ').title()}** "
            f"(Relevance: {item['similarity_score']:.3f}, Category: {item['category']}):\n"
            f"{item['content']}"
        )
    
    rag_context = "\n\n".join(context_sections)
    
    # Get prompt configuration
    config = prompt_manager.get_prompt_config("investment_analysis")
    
    # Enhanced system prompt with RAG context
    system_prompt_with_rag = f"""RETRIEVED INVESTMENT KNOWLEDGE & FRAMEWORKS:
{rag_context}

{config["system_prompt"]}

Use the retrieved investment knowledge above to inform your analysis and provide industry-standard insights with specific benchmarks and frameworks."""
    
    user_prompt = config["user_prompt_template"].format(content=content)
    model_config = config["model_config"]

    client = OpenAI(api_key=openai_api_key)

    response = client.chat.completions.create(
        model=model_config["model"],
        messages=[
            {"role": "system", "content": system_prompt_with_rag},
            {"role": "user", "content": user_prompt}
        ],
        max_tokens=model_config["max_tokens"],
        temperature=model_config["temperature"]
    )
    return response.choices[0].message.content

# 4. Main function to process a pitch deck file and save the Markdown report
def process_pitch_deck(file_path, openai_api_key):
    if file_path.lower().endswith(".pdf"):
        text = extract_text_from_pdf(file_path)
    elif file_path.lower().endswith(".pptx"):
        text = extract_text_from_pptx(file_path)
    else:
        raise ValueError("Unsupported file format. Please provide a PDF or PPTX file.")

    print("Extracted text length:", len(text))

    # Initialize systems
    try:
        print("Initializing RAG system...")
        rag_system = FAISSRAGSystem()
        
        print("Loading prompt configuration...")
        prompt_manager = PromptManager()
        
        print("Generating RAG-enhanced analysis...")
        markdown_report = analyze_pitch_deck_with_rag(text, openai_api_key, prompt_manager, rag_system)
        
    except (FileNotFoundError, ValueError) as e:
        print(f"Error during analysis: {e}")
        return

    output_markdown_path = file_path.rsplit(".", 1)[0] + "_rag_analysis.md"
    with open(output_markdown_path, "w", encoding="utf-8") as f:
        f.write(f"# Investment Analysis Report (RAG-Enhanced)\n")
        f.write(f"**Generated:** {time.strftime('%Y-%m-%d %H:%M:%S')}\n")
        f.write(f"**Source:** {os.path.basename(file_path)}\n\n")
        f.write(markdown_report)

    print(f"RAG-enhanced investment analysis saved to: {output_markdown_path}")

if __name__ == "__main__":
    # Check if API key is provided as argument or environment variable
    if len(sys.argv) == 3:
        file_path = sys.argv[1]
        openai_api_key = sys.argv[2]
    elif len(sys.argv) == 2:
        file_path = sys.argv[1]
        openai_api_key = os.getenv('OPENAI_API_KEY')
        if not openai_api_key:
            print("Error: OPENAI_API_KEY environment variable not set")
            print("Usage: python app.py <path_to_pitchdeck.pdf/pptx> [openai_api_key]")
            sys.exit(1)
    else:
        print("Usage: python app.py <path_to_pitchdeck.pdf/pptx> [openai_api_key]")
        print("Or set OPENAI_API_KEY environment variable")
        sys.exit(1)
        
    process_pitch_deck(file_path, openai_api_key)
